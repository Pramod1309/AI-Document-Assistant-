import os
import pdfkit
import subprocess
from celery import shared_task
from PIL import Image
from PyPDF2 import PdfWriter, PdfReader
import io
from django.conf import settings
import logging
from urllib.parse import quote
import multiprocessing
import concurrent.futures
import re
from docx import Document
import uuid
from pdf2docx import Converter
import comtypes.client
import openpyxl
from pptx import Presentation
import fitz
import pythoncom
from pathlib import Path

logger = logging.getLogger(__name__)

wkhtmltopdf_path = r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe'
config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path)

def parse_size_to_bytes(size_str):
    size_str = size_str.lower().replace(" ", "")
    match = re.match(r'(\d+\.?\d*)(kb|mb)', size_str)
    if match:
        value = float(match.group(1))
        unit = match.group(2)
        return int(value * (1024 if unit == 'kb' else 1024 * 1024))
    return None

def parse_resolution(resolution_str):
    match = re.match(r'(\d+)\s*[xX]\s*(\d+)', resolution_str)
    if match:
        return int(match.group(1)), int(match.group(2))
    return None, None

def parse_aspect_ratio(aspect_str):
    match = re.match(r'(\d+):(\d+)', aspect_str)
    if match:
        return int(match.group(1)), int(match.group(2))
    return None, None

def cleanup_files(*paths):
    for path in paths:
        try:
            if os.path.exists(path):
                os.remove(path)
        except Exception as e:
            logger.error(f"Failed to delete {path}: {str(e)}")

def resize_image(image_path, output_path, params):
    try:
        with Image.open(image_path) as img:
            original_width, original_height = img.size
            new_width, new_height = original_width, original_height
            quality = 95

            if params.get('size'):
                target_bytes = parse_size_to_bytes(params['size'])
                if target_bytes:
                    current_size = os.path.getsize(image_path)
                    quality = max(10, min(95, int(95 * (target_bytes / current_size))))
            elif params.get('height') and params.get('width'):
                new_width, new_height = int(params['width']), int(params['height'])
            elif params.get('aspect'):
                aspect_w, aspect_h = parse_aspect_ratio(params['aspect'])
                if aspect_w and aspect_h:
                    aspect_ratio = aspect_w / aspect_h
                    new_height = int(original_width / aspect_ratio)
                    if new_height > original_height:
                        new_height = original_height
                        new_width = int(original_height * aspect_ratio)
                    else:
                        new_width = int(new_height * aspect_ratio)

            img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            img_resized.save(output_path, optimize=True, quality=quality)
            return output_path
    except Exception as e:
        logger.error(f"Error resizing image {image_path}: {str(e)}")
        raise

def convert_docx_to_pdf(input_path, output_path):
    try:
        logger.info(f"Converting DOCX to PDF: {input_path} -> {output_path}")
        doc = Document(input_path)
        logger.info(f"Loaded DOCX with {len(doc.paragraphs)} paragraphs")
        html_content = "<html><body><h1>Converted Document</h1>"
        for para in doc.paragraphs:
            html_content += f"<p>{para.text}</p>"
        html_content += "</body></html>"
        
        temp_html_path = os.path.join(settings.MEDIA_ROOT, "temp", f"temp_{uuid.uuid4()}.html")
        os.makedirs(os.path.dirname(temp_html_path), exist_ok=True)
        logger.info(f"Writing temporary HTML to {temp_html_path}")
        with open(temp_html_path, "w", encoding="utf-8") as html_file:
            html_file.write(html_content)
        
        pdfkit.from_file(
            temp_html_path,
            output_path,
            configuration=config,
            options={
                "load-error-handling": "ignore",
                "enable-local-file-access": None,
                "quiet": "",
                "--dpi": "300",
                "--image-quality": "100",
            }
        )
        logger.info(f"Generated PDF at {output_path}, size: {os.path.getsize(output_path)} bytes")
        cleanup_files(temp_html_path)
    except Exception as e:
        logger.error(f"Error converting {input_path} to PDF: {str(e)}", exc_info=True)
        cleanup_files(temp_html_path)
        raise

def convert_to_pdf(input_path, output_path):
    if input_path.lower().endswith('.docx'):
        convert_docx_to_pdf(input_path, output_path)
    else:
        pdfkit.from_file(
            input_path,
            output_path,
            configuration=config,
            options={
                "load-error-handling": "ignore",
                "enable-local-file-access": None,
                "quiet": "",
                "--dpi": "300",
                "--image-quality": "100",
            }
        )

def compress_with_ghostscript(input_pdf_path, output_pdf_path, target_size_bytes):
    try:
        initial_size = os.path.getsize(input_pdf_path)
        quality_factor = 0.9

        while True:
            command = [
                "C:\\Program Files\\gs\\gs10.05.0\\bin\\gswin64c",
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                "-dPDFSETTINGS=/screen",
                "-dColorImageDownsampleType=/Bicubic",
                "-dColorImageResolution=72",
                "-dGrayImageDownsampleType=/Bicubic",
                "-dGrayImageResolution=72",
                "-dMonoImageDownsampleType=/Subsample",
                "-dMonoImageResolution=72",
                "-dDownsampleColorImages=true",
                "-dDownsampleGrayImages=true",
                "-dDownsampleMonoImages=true",
                "-dDetectDuplicateImages=true",
                "-dAutoFilterColorImages=false",
                "-dAutoFilterGrayImages=false",
                "-dQFactor={}".format(quality_factor * 1.0),
                "-dColorImageQuality={}".format(int(quality_factor * 100)),
                "-dGrayImageQuality={}".format(int(quality_factor * 100)),
                "-dNOPAUSE",
                "-dQUIET",
                "-dBATCH",
                f"-sOutputFile={output_pdf_path}",
                input_pdf_path,
            ]
            subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)

            current_size = os.path.getsize(output_pdf_path)
            logger.info(f"Compressed size: {current_size} bytes, Target: {target_size_bytes} bytes")

            if current_size <= target_size_bytes or quality_factor <= 0.1:
                break

            quality_factor -= 0.1

        with open(output_pdf_path, 'a+') as f:
            f.flush()
            os.fsync(f.fileno())
            
        if os.path.getsize(output_pdf_path) == 0:
            raise ValueError("Compressed PDF is empty")
            
        return output_pdf_path
    except subprocess.CalledProcessError as e:
        logger.error(f"Ghostscript error: {e.stderr.decode()}")
        raise
    except Exception as e:
        logger.error(f"Error in compress_with_ghostscript: {str(e)}")
        raise

@shared_task(bind=True, max_retries=3)
def convert_and_compress_images_to_pdf(self, image_paths, output_pdf_path, compressed_pdf_path, desired_size_str):
    try:
        temp_dir = os.path.join(settings.MEDIA_ROOT, "temp")
        os.makedirs(temp_dir, exist_ok=True)

        desired_size_bytes = parse_size_to_bytes(desired_size_str)
        if not desired_size_bytes:
            raise ValueError("Invalid desired size format.")

        resized_image_paths = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=multiprocessing.cpu_count() * 2) as executor:
            futures = []
            for i, image_path in enumerate(image_paths):
                resized_path = os.path.join(temp_dir, f"resized_{i}_{os.path.basename(image_path)}")
                futures.append(executor.submit(resize_image, image_path, resized_path, {'size': desired_size_str}))
            resized_image_paths = [future.result() for future in concurrent.futures.as_completed(futures)]

        html_content = "<html><body>"
        for resized_path in resized_image_paths:
            abs_path = os.path.abspath(resized_path)
            image_url = f"file:///{quote(abs_path.replace(os.sep, '/'))}"
            html_content += f'<img src="{image_url}" style="max-width: 100%; page-break-after: always;"><br>'
        html_content += "</body></html>"

        temp_html_path = os.path.join(temp_dir, f"temp_{self.request.id}.html")
        with open(temp_html_path, "w", encoding="utf-8") as html_file:
            html_file.write(html_content)

        pdfkit.from_file(
            temp_html_path,
            output_pdf_path,
            configuration=config,
            options={
                "load-error-handling": "ignore",
                "enable-local-file-access": None,
                "quiet": "",
                "--dpi": "300",
                "--image-quality": "100",
            }
        )

        initial_size = os.path.getsize(output_pdf_path)
        logger.info(f"Initial PDF size: {initial_size} bytes, Desired size: {desired_size_bytes} bytes")
        
        if initial_size > desired_size_bytes:
            compress_with_ghostscript(output_pdf_path, compressed_pdf_path, desired_size_bytes)
        else:
            os.rename(output_pdf_path, compressed_pdf_path)

        for path in [output_pdf_path, compressed_pdf_path]:
            if os.path.exists(path):
                with open(path, 'rb+') as f:
                    f.flush()
                    os.fsync(f.fileno())
                if os.path.getsize(path) == 0:
                    raise ValueError(f"Output file {path} is empty")
                if path.endswith('.pdf') and not verify_pdf_integrity(path):
                    raise ValueError(f"Generated PDF {path} fails integrity check")

        cleanup_files(temp_html_path, *resized_image_paths, *image_paths)

        return {
            "converted": output_pdf_path,
            "compressed": compressed_pdf_path,
            "file_size": os.path.getsize(compressed_pdf_path)
        }
    except Exception as e:
        logger.error(f"Error in convert_and_compress_images_to_pdf: {str(e)}")
        cleanup_files(*image_paths)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def convert_parallel_operations(self, first_input_path, second_input_path, first_output_path, second_output_path, first_op, second_op, params=None):
    try:
        # Normalize and validate input paths
        first_input_path = str(Path(first_input_path).resolve())
        second_input_path = str(Path(second_input_path).resolve())
        first_output_path = str(Path(first_output_path).resolve())
        second_output_path = str(Path(second_output_path).resolve())

        # Validate input files
        for path in [first_input_path, second_input_path]:
            if not os.path.exists(path):
                logger.error(f"Input file not found: {path}")
                raise FileNotFoundError(f"Input file not found: {path}")

        logger.info(f"Starting parallel operations: {first_op} on {first_input_path}, {second_op} on {second_input_path}")

        # Map operation names (handle "resize" as "resize_image")
        operation_mapping = {
            "convert_to_pdf": convert_to_pdf,
            "resize_image": resize_image,
            "resize": resize_image  # Add alias for "resize"
        }

        # Validate operations
        if first_op not in operation_mapping:
            raise ValueError(f"Unsupported first operation: {first_op}")
        if second_op not in operation_mapping:
            raise ValueError(f"Unsupported second operation: {second_op}")

        with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
            # Submit tasks with appropriate arguments
            if first_op == "resize_image" or first_op == "resize":
                future1 = executor.submit(operation_mapping[first_op], first_input_path, first_output_path, params)
            else:
                future1 = executor.submit(operation_mapping[first_op], first_input_path, first_output_path)

            if second_op == "resize_image" or second_op == "resize":
                future2 = executor.submit(operation_mapping[second_op], second_input_path, second_output_path, params)
            else:
                future2 = executor.submit(operation_mapping[second_op], second_input_path, second_output_path)

            # Wait for results
            first_result = future1.result()
            second_result = future2.result()

        # Verify outputs
        for output_path in [first_output_path, second_output_path]:
            with open(output_path, 'rb+') as f:
                f.flush()
                os.fsync(f.fileno())
            if os.path.getsize(output_path) == 0:
                raise ValueError(f"Output file {output_path} is empty")
            if output_path.endswith('.pdf') and not verify_pdf_integrity(output_path):
                raise ValueError(f"Generated PDF {output_path} fails integrity check")

        # Clean up input files
        cleanup_files(first_input_path, second_input_path)

        # Return result
        return {
            "first_output": first_result,
            "second_output": second_result,
            "first_size": os.path.getsize(first_result),
            "second_size": os.path.getsize(second_result)
        }
    except Exception as e:
        logger.error(f"Error in convert_parallel_operations: {str(e)}")
        try:
            self.retry(exc=e, countdown=5)
        except self.MaxRetriesExceededError:
            cleanup_files(first_input_path, second_input_path)  # Cleanup only after max retries
            return {
                "first_output": None,
                "second_output": None,
                "first_size": 0,
                "second_size": 0,
                "error": str(e)
           }
        
@shared_task(bind=True, max_retries=3)
def images_to_pdf(self, image_paths, output_path):
    try:
        temp_dir = os.path.join(settings.MEDIA_ROOT, "temp")
        os.makedirs(temp_dir, exist_ok=True)

        html_content = "<html><body>"
        for image_path in image_paths:
            abs_path = os.path.abspath(image_path)
            image_url = f"file:///{quote(abs_path.replace(os.sep, '/'))}"
            html_content += f'<img src="{image_url}" style="max-width: 100%; page-break-after: always;"><br>'
        html_content += "</body></html>"

        temp_html_path = os.path.join(temp_dir, f"temp_{self.request.id}.html")
        with open(temp_html_path, "w", encoding="utf-8") as html_file:
            html_file.write(html_content)

        pdfkit.from_file(
            temp_html_path,
            output_path,
            configuration=config,
            options={
                "load-error-handling": "ignore",
                "enable-local-file-access": None,
                "quiet": "",
                "--dpi": "300",
                "--image-quality": "100",
            }
        )

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")
        if not verify_pdf_integrity(output_path):
            raise ValueError(f"Generated PDF {output_path} fails integrity check")

        cleanup_files(temp_html_path, *image_paths)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in images_to_pdf: {str(e)}")
        cleanup_files(*image_paths)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def compress_pdf(self, input_path, output_path, desired_size_str):
    try:
        desired_size_bytes = parse_size_to_bytes(desired_size_str)
        if not desired_size_bytes:
            raise ValueError("Invalid desired size format.")

        compress_with_ghostscript(input_path, output_path, desired_size_bytes)

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")
        if not verify_pdf_integrity(output_path):
            raise ValueError(f"Generated PDF {output_path} fails integrity check")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in compress_pdf: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def word_to_pdf(self, input_path, output_path):
    try:
        convert_docx_to_pdf(input_path, output_path)

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")
        if not verify_pdf_integrity(output_path):
            raise ValueError(f"Generated PDF {output_path} fails integrity check")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in word_to_pdf: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def pdf_to_word(self, input_path, output_path):
    try:
        if not verify_pdf_integrity(input_path):
            raise ValueError(f"Input PDF {input_path} is invalid")

        cv = Converter(input_path)
        cv.convert(output_path)
        cv.close()

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in pdf_to_word: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def ppt_to_pdf(self, input_path, output_path):
    try:
        pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        presentation = powerpoint.Presentations.Open(os.path.abspath(input_path))
        presentation.SaveAs(os.path.abspath(output_path), 32)  # 32 = ppSaveAsPDF
        presentation.Close()
        powerpoint.Quit()

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")
        if not verify_pdf_integrity(output_path):
            raise ValueError(f"Generated PDF {output_path} fails integrity check")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in ppt_to_pdf: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)
    finally:
        try:
            pythoncom.CoUninitialize()
        except:
            logger.warning("CoUninitialize failed, possibly already uninitialized")

@shared_task(bind=True, max_retries=3)
def excel_to_pdf(self, input_path, output_path):
    try:
        pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
        excel = comtypes.client.CreateObject("Excel.Application")
        excel.Visible = False
        workbook = excel.Workbooks.Open(os.path.abspath(input_path))
        workbook.ExportAsFixedFormat(0, os.path.abspath(output_path))  # 0 = xlTypePDF
        workbook.Close()
        excel.Quit()

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")
        if not verify_pdf_integrity(output_path):
            raise ValueError(f"Generated PDF {output_path} fails integrity check")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in excel_to_pdf: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)
    finally:
        try:
            pythoncom.CoUninitialize()
        except:
            logger.warning("CoUninitialize failed, possibly already uninitialized")

@shared_task(bind=True, max_retries=3)
def pdf_to_excel(self, input_path, output_path):
    try:
        if not verify_pdf_integrity(input_path):
            raise ValueError(f"Input PDF {input_path} is invalid")

        doc = fitz.open(input_path)
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = "Sheet1"

        row = 1
        for page in doc:
            text = page.get_text("text")
            if not text.strip():
                logger.warning(f"Page {page.number} has no extractable text")
                sheet.cell(row=row, column=1).value = "No text extracted"
                row += 1
                continue
            lines = text.split('\n')
            for line in lines:
                sheet.cell(row=row, column=1).value = line
                row += 1

        if row == 1:
            logger.warning(f"No text extracted from {input_path}")
            sheet.cell(row=1, column=1).value = "No text extracted"

        workbook.save(output_path)

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in pdf_to_excel: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def pdf_to_ppt(self, input_path, output_path):
    try:
        if not verify_pdf_integrity(input_path):
            raise ValueError(f"Input PDF {input_path} is invalid")

        doc = fitz.open(input_path)
        prs = Presentation()

        for page in doc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            text = page.get_text("text")
            if not text.strip():
                logger.warning(f"Page {page.number} has no extractable text")
                text = "No text extracted"
            tx_box = slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
            tf = tx_box.text_frame
            tf.text = text

        prs.save(output_path)

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in pdf_to_ppt: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def convert_image_format(self, input_path, output_path, format):
    try:
        with Image.open(input_path) as img:
            if format.upper() == 'JPEG':
                img = img.convert('RGB')  # JPEG doesn't support RGBA
            img.save(output_path, format.upper())

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in convert_image_format: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

@shared_task(bind=True, max_retries=3)
def resize_image_task(self, input_path, output_path, params):
    try:
        resize_image(input_path, output_path, params)

        if os.path.getsize(output_path) == 0:
            raise ValueError(f"Output file {output_path} is empty")

        cleanup_files(input_path)

        return {"output": output_path}
    except Exception as e:
        logger.error(f"Error in resize_image_task: {str(e)}")
        cleanup_files(input_path)
        self.retry(exc=e, countdown=5)

def cleanup_files(*file_paths):
    for file_path in file_paths:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                logger.debug(f"Cleaned up file: {file_path}")
        except Exception as e:
            logger.error(f"Failed to clean up file {file_path}: {str(e)}")

def verify_pdf_integrity(pdf_path):
    try:
        with open(pdf_path, 'rb') as f:
            header = f.read(10)
            if not header.startswith(b'%PDF-'):
                logger.error(f"Invalid PDF header in {pdf_path}: {header}")
                return False
            reader = PdfReader(pdf_path)
            # Allow PDFs with zero pages for edge cases
            return True
    except Exception as e:
        logger.error(f"PDF verification failed for {pdf_path}: {str(e)}")
        return False